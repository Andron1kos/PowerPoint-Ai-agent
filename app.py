from flask import Flask, request, send_file
from pptx import Presentation
import uuid

app = Flask(__name__)

@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    data = request.json
    slides = data.get("slides", [])

    prs = Presentation()

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        title.text = slide_data.get("title", "")

        bullets = slide_data.get("bullets", [])
        tf = body.text_frame

        if bullets:
            tf.text = bullets[0]

            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1

    filename = f"{uuid.uuid4()}.pptx"
    prs.save(filename)

    return send_file(filename, as_attachment=True)

import os

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port)